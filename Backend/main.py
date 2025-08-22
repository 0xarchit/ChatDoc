import os
import time
import uuid
from io import BytesIO
import warnings
from contextlib import asynccontextmanager
from dotenv import load_dotenv

from fastapi import FastAPI, UploadFile, HTTPException, Query, Form
from fastapi.middleware.cors import CORSMiddleware
from typing import Optional
from pydantic import BaseModel

from langchain_mistralai import MistralAIEmbeddings
from langchain_openai import ChatOpenAI
from langchain_milvus import Milvus
from langchain_core.documents import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import create_retrieval_chain
from langchain_core.prompts import ChatPromptTemplate

import PyPDF2
import pandas as pd
from pptx import Presentation
from docx import Document as DocxDocument

load_dotenv()

@asynccontextmanager
async def lifespan(app: FastAPI):
    _init_vector_store()
    yield

# app = FastAPI(lifespan=lifespan,docs_url=None, redoc_url=None)
app = FastAPI(lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

MISTRAL_API_KEY = os.getenv("MISTRAL_API_KEY", "")
ZILLIZ_URI = os.getenv("ZILLIZ_URI", "")
ZILLIZ_TOKEN = os.getenv("ZILLIZ_TOKEN", "")
HF_TOKEN = os.getenv("HF_TOKEN", "")
COLLECTION_NAME = os.getenv("COLLECTION_NAME", "")
PASSWORD = os.getenv("PASSWORD", "")

if HF_TOKEN:
    os.environ.setdefault("HUGGINGFACEHUB_API_TOKEN", HF_TOKEN)
    os.environ.setdefault("HUGGING_FACE_HUB_TOKEN", HF_TOKEN)
    os.environ.setdefault("HUGGINGFACE_HUB_TOKEN", HF_TOKEN)



warnings.filterwarnings(
    "ignore",
    message=r".*Could not download mistral tokenizer.*",
    category=UserWarning,
    module=r"langchain_mistralai\.embeddings",
)

_VECTOR_STORE = None
_EMBEDDINGS = None

def _resolve_settings(
    mistral_key: Optional[str] = None,
    zilliz_uri: Optional[str] = None,
    zilliz_token: Optional[str] = None,
    collection_name: Optional[str] = None,
):
    """Resolve settings with priority: overrides > env vars > defaults."""
    return (
        mistral_key or os.getenv("MISTRAL_API_KEY") or MISTRAL_API_KEY,
        zilliz_uri or os.getenv("ZILLIZ_URI") or ZILLIZ_URI,
        zilliz_token or os.getenv("ZILLIZ_TOKEN") or ZILLIZ_TOKEN,
        collection_name or COLLECTION_NAME,
    )


def _build_vector_store(
    mistral_key: Optional[str] = None,
    zilliz_uri: Optional[str] = None,
    zilliz_token: Optional[str] = None,
    collection_name: Optional[str] = None,
):
    mk, uri, tok, coll = _resolve_settings(mistral_key, zilliz_uri, zilliz_token, collection_name)
    embeddings = MistralAIEmbeddings(
        model="mistral-embed",
        api_key=mk,
    )
    return Milvus(
        embedding_function=embeddings,
        connection_args={
            "uri": uri,
            "token": tok,
        },
        collection_name=coll,
    drop_old=False,
    auto_id=True,
    )


def _init_vector_store():
    global _VECTOR_STORE, _EMBEDDINGS
    mk, uri, tok, coll = _resolve_settings()
    _EMBEDDINGS = MistralAIEmbeddings(
        model="mistral-embed",
        api_key=mk,
    )
    _VECTOR_STORE = Milvus(
        embedding_function=_EMBEDDINGS,
        connection_args={
            "uri": uri,
            "token": tok,
        },
        collection_name=coll,
    drop_old=False,
    auto_id=True,
    )

def get_vector_store():
    if _VECTOR_STORE is None:
        _init_vector_store()
    return _VECTOR_STORE


def extract_text_from_file(content: bytes, filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    text = ""
    if ext == ".pdf":
        reader = PyPDF2.PdfReader(BytesIO(content))
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    elif ext == ".txt":
        for enc in ("utf-8", "utf-16", "utf-16le", "utf-16be", "cp1252"):
            try:
                text = content.decode(enc)
                break
            except UnicodeDecodeError:
                continue
        else:
            text = content.decode("utf-8", errors="ignore")
    elif ext == ".csv":
        df = pd.read_csv(BytesIO(content))
        text = df.to_string()
    elif ext == ".xlsx":
        df = pd.read_excel(BytesIO(content))
        text = df.to_string()
    elif ext in [".ppt", ".pptx"]:
        prs = Presentation(BytesIO(content))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif ext in [".doc", ".docx"]:
        doc = DocxDocument(BytesIO(content))
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    else:
        raise ValueError("Unsupported file type")
    return text

MAX_FILE_SIZE = 20 * 1024 * 1024 

@app.post("/upload")
async def upload_file(
    file: UploadFile,
    mistral_api_key: Optional[str] = Form(default=None),
    zilliz_uri: Optional[str] = Form(default=None),
    zilliz_token: Optional[str] = Form(default=None),
    collection_name: Optional[str] = Form(default=None),
):
    if not file:
        raise HTTPException(status_code=400, detail="No file uploaded")
    
    upload_id = str(uuid.uuid4())
    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="File too large. Max size is 20MB.")
    filename = file.filename
    
    try:
        text = extract_text_from_file(content, filename)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error extracting text: {str(e)}")
    if not text or not text.strip():
        raise HTTPException(status_code=400, detail="text extraction failed")
    
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    chunks = splitter.split_text(text)
    if not chunks or all(not c.strip() for c in chunks):
        raise HTTPException(status_code=400, detail="text extraction failed")
    
    docs = [
        Document(
            page_content=chunk,
            metadata={
                "upload_id": upload_id,
                "source": filename,
                "upload_time": time.time()
            }
        ) for chunk in chunks
    ]
    
    if any([mistral_api_key, zilliz_uri, zilliz_token, collection_name]):
        vector_store = _build_vector_store(mistral_api_key, zilliz_uri, zilliz_token, collection_name)
    else:
        vector_store = get_vector_store()
    vector_store.add_documents(docs)
    
    return {"upload_id": upload_id}

class QueryRequest(BaseModel):
    question: str
    upload_id: str
    mistral_api_key: Optional[str] = None
    zilliz_uri: Optional[str] = None
    zilliz_token: Optional[str] = None
    collection_name: Optional[str] = None

@app.post("/query")
def query_rag(request: QueryRequest):
    if any([request.mistral_api_key, request.zilliz_uri, request.zilliz_token, request.collection_name]):
        vector_store = _build_vector_store(
            request.mistral_api_key, request.zilliz_uri, request.zilliz_token, request.collection_name
        )
        mk, _, _, _ = _resolve_settings(request.mistral_api_key, request.zilliz_uri, request.zilliz_token, request.collection_name)
    else:
        vector_store = get_vector_store()
        mk, _, _, _ = _resolve_settings()

    llm = ChatOpenAI(
        model="gpt-5-nano",
        api_key="not-req",
        base_url="https://text.pollinations.ai/openai",
    )
    
    system_prompt = (
        "You are an assistant for question-answering tasks. "
        "Use the following pieces of retrieved context to answer "
        "the question. If you don't know the answer, say that you "
        "don't know. Use three sentences maximum and keep the "
        "answer concise."
        "If Anything out of data from file is asked then say "
        "It is not related to uploaded document. Please ask Something Valid"
        "\n\n"
        "{context}"
    )
    
    prompt = ChatPromptTemplate.from_messages(
        [
            ("system", system_prompt),
            ("human", "{input}"),
        ]
    )
    
    retriever = vector_store.as_retriever(
        search_kwargs={"expr": f'upload_id == "{request.upload_id}"', "k": 4}
    )
    
    chain = create_retrieval_chain(retriever, prompt | llm)
    
    try:
        response = chain.invoke({"input": request.question})
        return {"answer": response["answer"]}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating response: {str(e)}")

@app.delete("/delete/{upload_id}")
def delete_upload(
    upload_id: str,
    mistral_api_key: Optional[str] = Query(default=None),
    zilliz_uri: Optional[str] = Query(default=None),
    zilliz_token: Optional[str] = Query(default=None),
    collection_name: Optional[str] = Query(default=None),
):
    if any([mistral_api_key, zilliz_uri, zilliz_token, collection_name]):
        vector_store = _build_vector_store(mistral_api_key, zilliz_uri, zilliz_token, collection_name)
    else:
        vector_store = get_vector_store()
    vector_store.delete(expr=f'upload_id == "{upload_id}"')
    return {"status": "deleted"}

@app.get("/deleteall")
def delete_all(
    password: Optional[str] = Query(None, description="Admin password required when using native DB"),
    mistral_api_key: Optional[str] = Query(default=None),
    zilliz_uri: Optional[str] = Query(default=None),
    zilliz_token: Optional[str] = Query(default=None),
    collection_name: Optional[str] = Query(default=None),
):
    using_overrides = any([mistral_api_key, zilliz_uri, zilliz_token, collection_name])
    if using_overrides:
        vector_store = _build_vector_store(mistral_api_key, zilliz_uri, zilliz_token, collection_name)
    else:
        if password != PASSWORD:
            raise HTTPException(status_code=403, detail="Forbidden: invalid password")
        vector_store = get_vector_store()
    vector_store.delete(expr='upload_id like "%"')
    return {"status": "all_deleted"}

@app.get("/active")
def active():
    return {"message": "API is Active"}

# uvicorn main:app --reload